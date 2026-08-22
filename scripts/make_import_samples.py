"""読み込みを試すためのファイルを作る（F-2 / F-3 / ADR-0027）。

**同じ内容を 4 つの形で出す。** 見出し・段落・箇条書き・表・図形を 1 枚に
入れてあるので、どの経路で取り込んでも「何が落ちるか」を並べて見られる。

    uv run python scripts/make_import_samples.py [出力先]

出るもの（既定は `samples/`）:

| ファイル              | 通る道                                    |
| --------------------- | ----------------------------------------- |
| `会議メモ.pdf`        | 文字入りの PDF（QtPdf がそのまま取り出す）|
| `会議メモ-スキャン.pdf` | 絵だけの PDF（読み取りに回る）          |
| `会議メモ.png`        | 画像（読み取りに回る）                    |
| `会議メモ.jpg`        | 同上（写真に近い圧縮）                    |
| `会議メモ.pptx`       | PowerPoint（`pptx_import`）               |
| `会議メモ-図つき.pdf` | 文字と図が同じページ（図は添付になる）    |
| `会議メモ-混在.pdf`   | 1 枚目は文字・2 枚目は絵（ページごとに判断）|

**作ったものは `samples/` に入れてある**（読み込みを試すのに毎回作らなくて
よいように。ユーザー決定 2026-08-23）。中身を変えたら作り直して入れ替える。
"""

import sys
from pathlib import Path

from PySide6.QtCore import QPointF, QRectF, Qt
from PySide6.QtGui import (
    QBrush,
    QColor,
    QFont,
    QImage,
    QPageSize,
    QPainter,
    QPdfWriter,
    QPen,
    QPolygonF,
)
from PySide6.QtWidgets import QApplication

TITLE = "会議メモ 2026-08-22"

PARAGRAPH = (
    "来週の予算について話し合った。営業部からは増額の要望が出ているが、"
    "根拠となる数字はまだ揃っていない。"
)

BULLETS = [
    "前年比の資料を用意する",
    "判断は経営会議で行う",
    "日程はまだ決まっていない",
]

TABLE_HEADER = ["項目", "金額", "担当"]
TABLE_ROWS = [
    ["会場費", "1,250,000 円", "野村"],
    ["印刷費", "84,000 円", "田中"],
]

FLOW = ["申請", "承認", "発注"]
"""図形（箱と矢印）。**絵として置く**ので、読み取りでは箱の中の字だけが残る。"""

# 1 枚の絵の大きさ（A4 に近い比）。読み取りに回すので、小さすぎると読めない
PAGE_WIDTH = 1240
PAGE_HEIGHT = 1754


def draw_page(painter: QPainter, width: int, height: int) -> None:
    """1 枚ぶんを描く。**PDF にも画像にも同じ手で描く**（中身を揃えるため）。"""
    painter.setRenderHint(QPainter.RenderHint.Antialiasing)
    painter.fillRect(QRectF(0, 0, width, height), QColor("white"))
    painter.setPen(QColor("black"))

    scale = width / PAGE_WIDTH
    margin = 90 * scale
    y = margin

    def text(value: str, size: float, *, bold: bool = False, indent: float = 0.0) -> None:
        nonlocal y
        font = QFont("Hiragino Sans", int(size * scale))
        font.setBold(bold)
        painter.setFont(font)
        box = QRectF(margin + indent, y, width - margin * 2 - indent, 200 * scale)
        flags = int(Qt.AlignmentFlag.AlignTop | Qt.TextFlag.TextWordWrap)
        painter.drawText(box, flags, value)
        y += painter.boundingRect(box, flags, value).height() + 14 * scale

    text(TITLE, 30, bold=True)
    y += 10 * scale
    text(PARAGRAPH, 19)
    y += 16 * scale

    text("決まったこと", 23, bold=True)
    for item in BULLETS:
        text(f"・{item}", 19, indent=20 * scale)
    y += 16 * scale

    text("費用", 23, bold=True)
    y = _draw_table(painter, margin, y, width - margin * 2, scale)
    y += 30 * scale

    text("進め方", 23, bold=True)
    _draw_flow(painter, margin, y, scale)


def _draw_table(painter: QPainter, left: float, top: float, width: float, scale: float) -> float:
    """罫線付きの表。**絵として引く**ので、読み取りでは線が消えて字だけ残る。"""
    columns = [width * 0.45, width * 0.35, width * 0.20]
    row_height = 46 * scale
    font = QFont("Hiragino Sans", int(18 * scale))
    painter.setPen(QPen(QColor("#333333"), 1.2 * scale))

    y = top
    for index, row in enumerate([TABLE_HEADER, *TABLE_ROWS]):
        x = left
        font.setBold(index == 0)
        painter.setFont(font)
        for column, value in zip(columns, row, strict=True):
            cell = QRectF(x, y, column, row_height)
            painter.drawRect(cell)
            painter.drawText(
                cell.adjusted(10 * scale, 0, -10 * scale, 0),
                int(Qt.AlignmentFlag.AlignVCenter | Qt.AlignmentFlag.AlignLeft),
                value,
            )
            x += column
        y += row_height
    return y


def _draw_flow(painter: QPainter, left: float, top: float, scale: float) -> None:
    """箱と矢印。**図形は読み取りでは字だけになる**ことを確かめるために置く。"""
    box_width = 190 * scale
    box_height = 74 * scale
    gap = 70 * scale
    font = QFont("Hiragino Sans", int(18 * scale))
    painter.setFont(font)

    x = left
    for index, label in enumerate(FLOW):
        box = QRectF(x, top, box_width, box_height)
        painter.setPen(QPen(QColor("#2C5AA0"), 2 * scale))
        painter.setBrush(QBrush(QColor("#EAF1FB")))
        painter.drawRoundedRect(box, 8 * scale, 8 * scale)
        painter.setPen(QColor("black"))
        painter.drawText(box, int(Qt.AlignmentFlag.AlignCenter), label)

        if index < len(FLOW) - 1:
            _draw_arrow(painter, box.right(), top + box_height / 2, gap, scale)
        x += box_width + gap
    painter.setBrush(Qt.BrushStyle.NoBrush)


def _draw_arrow(painter: QPainter, x: float, y: float, length: float, scale: float) -> None:
    painter.setPen(QPen(QColor("#2C5AA0"), 2 * scale))
    tip = QPointF(x + length * 0.75, y)
    painter.drawLine(QPointF(x + length * 0.1, y), tip)
    head = QPolygonF(
        [
            tip,
            QPointF(tip.x() - 12 * scale, y - 8 * scale),
            QPointF(tip.x() - 12 * scale, y + 8 * scale),
        ]
    )
    painter.setBrush(QBrush(QColor("#2C5AA0")))
    painter.drawPolygon(head)


def figure_image() -> QImage:
    """棒グラフ 1 枚（ADR-0027 追記）。**文字ではなく図**を置くために作る。

    文字の写った絵を置くと「読み取られるのか、添付になるのか」が
    分からなくなるので、**読ませる意味の無い図**にしてある。
    """
    width, height = 520, 300
    image = QImage(width, height, QImage.Format.Format_RGB32)
    image.fill(QColor("white"))
    painter = QPainter(image)
    painter.setRenderHint(QPainter.RenderHint.Antialiasing)
    painter.setFont(QFont("Hiragino Sans", 13))
    painter.setPen(QColor("#333333"))
    painter.drawText(QRectF(16, 10, width - 32, 30), int(Qt.AlignmentFlag.AlignLeft), "費用の内訳")
    painter.drawLine(QPointF(60, height - 50), QPointF(width - 30, height - 50))

    bars = [("会場費", 0.92), ("印刷費", 0.28), ("予備", 0.15)]
    for index, (label, ratio) in enumerate(bars):
        left = 90 + index * 130
        bar_height = (height - 110) * ratio
        painter.setBrush(QBrush(QColor("#2C5AA0")))
        painter.setPen(Qt.PenStyle.NoPen)
        painter.drawRect(QRectF(left, height - 50 - bar_height, 70, bar_height))
        painter.setPen(QColor("#333333"))
        painter.drawText(
            QRectF(left - 15, height - 44, 100, 24), int(Qt.AlignmentFlag.AlignLeft), label
        )
    painter.end()
    return image


def write_pdf_with_figure(path: Path) -> None:
    """文字と図が**同じページ**にある PDF（ADR-0027 追記）。

    読み取りには回らず（文字がある）、**図だけが添付として取り込まれる**。
    """
    writer = QPdfWriter(str(path))
    writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
    writer.setResolution(150)
    painter = QPainter(writer)
    view = painter.viewport()
    scale = view.width() / PAGE_WIDTH
    margin = 90 * scale

    painter.setFont(QFont("Hiragino Sans", int(26 * scale)))
    painter.drawText(QRectF(margin, margin, view.width() - margin * 2, 80 * scale), 0, TITLE)
    painter.setFont(QFont("Hiragino Sans", int(17 * scale)))
    painter.drawText(
        QRectF(margin, margin + 90 * scale, view.width() - margin * 2, 300 * scale),
        int(Qt.AlignmentFlag.AlignTop | Qt.TextFlag.TextWordWrap),
        PARAGRAPH + "\n\n内訳は下の図のとおり。",
    )
    figure = figure_image()
    top = margin + 320 * scale
    painter.drawImage(
        QRectF(margin, top, figure.width() * scale * 1.6, figure.height() * scale * 1.6), figure
    )
    painter.end()


def write_mixed_pdf(path: Path, image: QImage) -> None:
    """1 枚目は文字、2 枚目は絵だけ（ページごとの切り分けを見るため）。"""
    writer = QPdfWriter(str(path))
    writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
    writer.setResolution(150)
    painter = QPainter(writer)
    view = painter.viewport()
    scale = view.width() / PAGE_WIDTH
    margin = 90 * scale

    painter.setFont(QFont("Hiragino Sans", int(17 * scale)))
    painter.drawText(
        QRectF(margin, margin, view.width() - margin * 2, 400 * scale),
        int(Qt.AlignmentFlag.AlignTop | Qt.TextFlag.TextWordWrap),
        "1 枚目は文字として入っています。ここは読み取りに回りません。\n\n"
        "2 枚目は紙を取り込んだページで、そちらだけが読み取りに回ります。",
    )
    writer.newPage()
    painter.drawImage(
        QRectF(0, 0, view.width(), view.width() * image.height() / image.width()), image
    )
    painter.end()


def write_pdf(path: Path) -> None:
    """文字が入った PDF。**取り込みは読み取りに回らない**（速くて正確）。"""
    writer = QPdfWriter(str(path))
    writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
    writer.setResolution(150)
    painter = QPainter(writer)
    view = painter.viewport()
    draw_page(painter, view.width(), view.height())
    painter.end()


def write_image(path: Path, *, quality: int = -1) -> QImage:
    """画像 1 枚。読み取りの経路を試す。"""
    image = QImage(PAGE_WIDTH, PAGE_HEIGHT, QImage.Format.Format_RGB32)
    painter = QPainter(image)
    draw_page(painter, PAGE_WIDTH, PAGE_HEIGHT)
    painter.end()
    image.save(str(path), quality=quality)
    return image


def write_scanned_pdf(path: Path, image: QImage) -> None:
    """絵だけの PDF（紙を取り込んだもの）。**文字は入っていない。**"""
    writer = QPdfWriter(str(path))
    writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
    writer.setResolution(150)
    painter = QPainter(writer)
    view = painter.viewport()
    height = view.width() * image.height() / image.width()
    painter.drawImage(QRectF(0, 0, view.width(), height), image)
    painter.end()


def write_pptx(path: Path) -> None:
    """PowerPoint。表と図形を**構造のまま**持つので、取り込みの比較になる。"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm, Pt

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = TITLE
    body = slide.placeholders[1].text_frame
    body.text = PARAGRAPH
    for item in BULLETS:
        paragraph = body.add_paragraph()
        paragraph.text = item
        paragraph.level = 1

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    table_slide.shapes.title.text = "費用"
    shape = table_slide.shapes.add_table(
        len(TABLE_ROWS) + 1, len(TABLE_HEADER), Cm(2), Cm(4), Cm(22), Cm(4)
    )
    for column, value in enumerate(TABLE_HEADER):
        shape.table.cell(0, column).text = value
    for row, values in enumerate(TABLE_ROWS, start=1):
        for column, value in enumerate(values):
            shape.table.cell(row, column).text = value

    flow_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    flow_slide.shapes.title.text = "進め方"
    for index, label in enumerate(FLOW):
        box = flow_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2 + index * 7), Cm(6), Cm(5), Cm(2)
        )
        box.text_frame.text = label
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        if index < len(FLOW) - 1:
            flow_slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Cm(7 + index * 7), Cm(6.6), Cm(2), Cm(0.8)
            )
    presentation.save(str(path))


def main(argv: list[str]) -> int:
    target = Path(argv[1]).expanduser() if len(argv) > 1 else Path("samples")
    target.mkdir(parents=True, exist_ok=True)

    app = QApplication.instance() or QApplication([])
    assert app is not None

    write_pdf(target / "会議メモ.pdf")
    image = write_image(target / "会議メモ.png")
    write_image(target / "会議メモ.jpg", quality=85)
    write_scanned_pdf(target / "会議メモ-スキャン.pdf", image)
    write_pdf_with_figure(target / "会議メモ-図つき.pdf")
    write_mixed_pdf(target / "会議メモ-混在.pdf", image)
    write_pptx(target / "会議メモ.pptx")

    for path in sorted(target.iterdir()):
        print(f"{path}  {path.stat().st_size / 1024:.0f}KB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
