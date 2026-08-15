"""PowerPoint への書き出し（F-5）。

**ざっくり作って手で整える**前提。凝ったレイアウトは狙わない。

割り方は `core/slides.py`（F-4）が決めていて、ここは組み立てだけを持つ。
分けてあるので、規則を変えたいときに触る場所が 1 つで済む。

置き方はユーザーと決めた。**`#` は表紙、`##` ごとに 1 枚、画像は右側。**
画像があるスライドは本文を左半分に寄せる（画像と重ならないように）。

**書き出しは止めない。** 画像が見つからなくても、保管フォルダの外を
指していても、そこだけ飛ばしてファイルを作る。1 枚のリンク切れで
書き出せないほうが困る。
"""

import logging
from pathlib import Path

from hitofude.core.paths import resolve_reference
from hitofude.core.slides import Block, BlockKind, Deck, Slide, split
from hitofude.core.table import split_cells

logger = logging.getLogger(__name__)

# スライドの大きさ（16:9）。既定の 4:3 は今どき狭い
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

# 余白と本文の置き場所（インチ）
MARGIN = 0.6
TITLE_TOP = 0.5
TITLE_HEIGHT = 1.1
BODY_TOP = 1.8
BODY_HEIGHT = 5.0

# 画像があるときの本文の幅（全体に対する割合）。残りが画像の場所になる
BODY_RATIO_WITH_IMAGE = 0.52

# 文字の大きさ
TITLE_POINTS = 30
BODY_POINTS = 17
HEADING_POINTS = 19
CODE_POINTS = 13
TABLE_POINTS = 13

# PowerPoint の箇条書きは 0〜8 段
MAX_LEVEL = 8

# 本文の枠に入れる種類。コードと表は別の図形にする（枠に入らない）
_TEXT_KINDS = frozenset({BlockKind.PARAGRAPH, BlockKind.BULLET, BlockKind.HEADING})

# 既定のレイアウト番号（`python-pptx` の既定テンプレート）
LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6


def write_pptx(path: Path, text: str, *, base_path: Path | None = None) -> Path:
    """本文を `.pptx` にする。書いた先を返す。

    `base_path` は画像を探す起点（保管フォルダ）。渡さなければ画像は入らない。
    """
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)

    deck = split(text)
    if deck.title:
        _add_title_slide(presentation, deck)
    for slide in deck.slides:
        _add_slide(presentation, slide, base_path)

    presentation.save(str(path))
    return path


def _add_title_slide(presentation, deck: Deck) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[LAYOUT_TITLE])
    slide.shapes.title.text = deck.title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = deck.subtitle


def _add_slide(presentation, slide: Slide, base_path: Path | None) -> None:
    from pptx.util import Inches, Pt

    built = presentation.slides.add_slide(presentation.slide_layouts[LAYOUT_CONTENT])
    if built.shapes.title is not None:
        built.shapes.title.text = slide.title
        for paragraph in built.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(TITLE_POINTS)

    images = _resolve_images(slide, base_path)
    body_width = (SLIDE_WIDTH_IN - MARGIN * 2) * (BODY_RATIO_WITH_IMAGE if images else 1.0)

    text_blocks = [block for block in slide.blocks if block.kind in _TEXT_KINDS]
    other_blocks = [block for block in slide.blocks if block.kind not in _TEXT_KINDS]

    body = built.placeholders[1] if len(built.placeholders) > 1 else None
    top = BODY_TOP
    if body is not None and text_blocks:
        body.left, body.top = Inches(MARGIN), Inches(BODY_TOP)
        body.width, body.height = Inches(body_width), Inches(_body_height(text_blocks))
        _fill_body(body.text_frame, text_blocks)
        top = BODY_TOP + _body_height(text_blocks) + 0.2
    elif body is not None:
        # 空の枠は「テキストを入力」と出るだけで邪魔になる
        body._element.getparent().remove(body._element)

    for block in other_blocks:
        top = _add_block(built, block, top, body_width)

    for image in images:
        _add_image(built, image, body_width)

    if slide.notes:
        built.notes_slide.notes_text_frame.text = slide.notes


def _fill_body(frame, blocks: list[Block]) -> None:
    """本文の枠に段落を積む。

    **箇条書きは枠の機能に任せる。** `・` を文字として足すと、読み戻した
    ときに `- ・項目` になる（往復テストが見つけた）。段落の階層
    （`level`）を設定すれば PowerPoint が行頭記号を描く。

    段落と小見出しは行頭記号を消す（`buNone`）。取り込み（F-3）は
    **この印を小見出しの手掛かりにしている**ので、往復で形が保たれる。
    """
    from pptx.util import Pt

    frame.word_wrap = True
    for index, block in enumerate(blocks):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = block.text

        if block.kind is BlockKind.BULLET:
            paragraph.level = min(block.level, MAX_LEVEL)
            run.font.size = Pt(BODY_POINTS)
            continue

        _drop_bullet(paragraph)
        if block.kind is BlockKind.HEADING:
            run.font.size = Pt(HEADING_POINTS)
            run.font.bold = True
        else:
            run.font.size = Pt(BODY_POINTS)


def _drop_bullet(paragraph) -> None:
    """その段落の行頭記号を消す（`<a:buNone/>`）。"""
    from pptx.oxml.ns import qn

    properties = paragraph._p.get_or_add_pPr()
    properties.append(properties.makeelement(qn("a:buNone"), {}))


def _body_height(blocks: list[Block]) -> float:
    return sum(_text_height(block) for block in blocks)


def _add_block(slide, block: Block, top: float, width: float) -> float:
    """本文の枠に入らないもの（コードと表）を置いて、次の上端を返す。"""
    if block.kind is BlockKind.TABLE:
        return _add_table(slide, block, top, width)
    return _add_code(slide, block, top, width)


def _text_height(block: Block) -> float:
    """その行に要る高さ（インチ）。**ざっくりでよい**（枠は自動で伸びる）。"""
    return 0.45 if block.kind is BlockKind.HEADING else 0.38


def _add_code(slide, block: Block, top: float, width: float) -> float:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    lines = block.text.split("\n")
    height = max(0.4, 0.26 * len(lines) + 0.2)
    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF6)

    frame = box.text_frame
    frame.word_wrap = False
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        # **等幅で書く。** 取り込み（F-3）はこれをコードの手掛かりにする
        run.font.name = "Menlo"
        run.font.size = Pt(CODE_POINTS)
    return top + height + 0.15


def _add_table(slide, block: Block, top: float, width: float) -> float:
    from pptx.util import Inches, Pt

    rows = [_cells(line) for line in block.lines]
    rows = [row for row in rows if any(row)]
    if not rows:
        return top

    columns = max(len(row) for row in rows)
    height = 0.4 * len(rows)
    shape = slide.shapes.add_table(
        len(rows), columns, Inches(MARGIN), Inches(top), Inches(width), Inches(height)
    )
    table = shape.table
    for row_index, row in enumerate(rows):
        for column_index in range(columns):
            cell = table.cell(row_index, column_index)
            cell.text = row[column_index] if column_index < len(row) else ""
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(TABLE_POINTS)
    return top + height + 0.2


def _cells(line: str) -> list[str]:
    """表の 1 行をセルに割る。**前後のパイプは区切りではない。**

    外さずに割ると、両端に空のセルが 1 つずつ増える（往復テストが
    `|   | 担当 | 人数 |   |` を見つけた）。
    """
    body = line.strip()
    if body.startswith("|"):
        body = body[1:]
    if body.endswith("|"):
        body = body[:-1]
    return [cell.strip() for cell in split_cells(body)]


def _resolve_images(slide: Slide, base_path: Path | None) -> list[Path]:
    """本文の画像を実ファイルへ解決する。

    **保管フォルダの外は読まない**（`core/paths`）。本文は手で編集できるので、
    `../` で任意のファイルを埋め込ませない（書き出した資料は人に渡る）。
    """
    found: list[Path] = []
    for reference in slide.images:
        resolved = resolve_reference(base_path, reference)
        if resolved is None:
            logger.warning("画像を入れられなかった: %s", reference)
            continue
        found.append(resolved)
    return found


def _add_image(slide, image: Path, body_width: float) -> None:
    """本文の右側に置く（ユーザーと決めた並べ方）。"""
    from pptx.util import Inches

    left = MARGIN + body_width + 0.3
    width = SLIDE_WIDTH_IN - left - MARGIN
    if width <= 0:
        return
    try:
        slide.shapes.add_picture(str(image), Inches(left), Inches(BODY_TOP), width=Inches(width))
    except OSError:
        logger.warning("画像を読めなかった: %s", image)
