"""PowerPoint を読んで Markdown にする（F-3）。

**ざっくり読んで手で直す**前提。元のレイアウト（配色・段組み・位置）は
復元しない。復元できるのは**中身**だけで、これは形式の側に情報が
残っていないため。

手掛かりは実物（md2pptx が作った `.pptx`）を調べて決めた。

| 手掛かり | 使い道 |
| --- | --- |
| スライドのタイトル枠 | `##`（PDF の取り込みと同じ区切り） |
| `buNone`（行頭記号なし）で短い段落 | `###` |
| 文の終わりの記号で終わる段落 | 本文。それ以外は `- ` |
| 等幅フォント（Consolas / Menlo など） | コードブロック |
| 太字の run | `**強調**` |

**平文と第 1 階層の箇条書きは、形式の上では見分けが付かない。** PowerPoint
の本文枠は既定で全段落に行頭記号が付くためで、上の「文の終わりの記号」は
その埋め合わせ。外れることがあるが、**箇条書きが 1 つ本文になっても
目で見て直せる**。

`python-pptx` は PySide6 に依存しないが、取り込みの入口をここ 1 か所に
まとめたいので `editor/` に置く（PDF は QtPdf を使うため `core/` に置けない）。
"""

import logging
from collections.abc import Callable
from pathlib import Path

from hitofude.core.imported import (
    is_page_number,
    looks_like_heading,
    normalize_text,
)
from hitofude.core.table import format_table

logger = logging.getLogger(__name__)

SUFFIX = ".pptx"

# 等幅として扱うフォント名（小文字で部分一致）。コードブロックの手掛かり
MONO_FONTS = ("consolas", "menlo", "monaco", "courier", "mono", "source code", "sf mono")

# 箇条書き 1 段ぶんの字下げ。このアプリの既定に合わせる
INDENT = "    "

# 文の終わりに見える記号。ここで終わる段落は本文として扱う
_SENTENCE_END = "。．.！？!?"


def to_markdown(path: Path, *, save_image: Callable[[bytes, str], str | None] | None = None) -> str:
    """`.pptx` 1 つを Markdown にする。読めなければ空。

    `save_image` は画像を保存して**本文に挿す Markdown を返す**関数
    （`MainWindow.save_attachment`）。渡さなければ画像は飛ばす。

    **中身が無ければ空を返す。** 題名だけのノートを作らせない
    （画像だけの PDF で踏んだのと同じ穴）。
    """
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        presentation = Presentation(str(path))
    except (PackageNotFoundError, OSError, ValueError, KeyError):
        logger.warning("PowerPoint を読めなかった: %s", path)
        return ""

    parts: list[str] = []
    for slide in presentation.slides:
        parts.extend(_slide_blocks(slide, save_image))

    if not parts:
        logger.warning("文字を取り出せなかった: %s", path)
        return ""
    return f"# {path.stem}\n\n" + "\n\n".join(parts) + "\n"


def _slide_blocks(slide, save_image) -> list[str]:
    """1 枚ぶんのブロック。タイトル → 中身 → 発表者ノートの順。"""
    blocks: list[str] = []
    title = _title_of(slide)
    if title:
        blocks.append(f"## {title}")

    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        blocks.extend(_shape_blocks(shape, save_image))

    blocks.extend(_notes_blocks(slide))
    return blocks


def _title_of(slide) -> str:
    try:
        title = slide.shapes.title
    except (AttributeError, KeyError):
        return ""
    return normalize_text(title.text).strip() if title is not None else ""


def _shape_blocks(shape, save_image) -> list[str]:
    if shape.shape_type == 13:  # PICTURE
        return _picture_blocks(shape, save_image)
    if getattr(shape, "has_table", False) and shape.has_table:
        return _table_blocks(shape.table)
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        return _text_blocks(shape)
    return []


def _text_blocks(shape) -> list[str]:
    """テキスト枠。等幅ならコード、そうでなければ段落と箇条書き。"""
    text = normalize_text(shape.text_frame.text).strip()
    if not text or is_page_number(text):
        return []

    if _is_mono(shape):
        # **中は触らない。** 字下げも記号もコードの一部
        return [f"```\n{shape.text_frame.text.rstrip()}\n```"]

    blocks: list[str] = []
    bullets: list[str] = []

    def flush() -> None:
        if bullets:
            blocks.append("\n".join(bullets))
            bullets.clear()

    for paragraph in shape.text_frame.paragraphs:
        line = normalize_text(_paragraph_text(paragraph)).strip()
        if not line or is_page_number(line):
            continue

        if _is_sub_heading(paragraph, line):
            flush()
            blocks.append(f"### {line}")
            continue

        if _looks_like_body(line, paragraph):
            flush()
            blocks.append(line)
            continue

        bullets.append(f"{INDENT * paragraph.level}- {line}")

    flush()
    return blocks


def _paragraph_text(paragraph) -> str:
    """run の書式を記号に戻して繋ぐ。太字は `**`、等幅は `` ` ``。

    **等幅は段落の中に混ざる。** インラインコード（`` `AWS` ``）が
    そう書かれているので、枠ごとコードにせず記号で囲み直す。
    """
    pieces: list[str] = []
    for run in paragraph.runs:
        text = run.text
        if not text.strip():
            pieces.append(text)
            continue
        head = text[: len(text) - len(text.lstrip())]
        tail = text[len(text.rstrip()) :]
        body = text.strip()
        if _is_mono_run(run):
            body = f"`{body}`"
        if run.font.bold:
            body = f"**{body}**"
        pieces.append(f"{head}{body}{tail}")
    return "".join(pieces) or paragraph.text


def _is_sub_heading(paragraph, line: str) -> bool:
    """行頭記号を消してある短い段落は小見出し（md2pptx の `###`）。

    **短さも見る。** 他のツールが本文の段落にも `buNone` を付けることが
    あり、そのまま信じると本文が全部見出しになる。
    """
    return _has_bullet_none(paragraph) and looks_like_heading(line)


def _has_bullet_none(paragraph) -> bool:
    from pptx.oxml.ns import qn

    properties = paragraph._pPr
    return properties is not None and properties.find(qn("a:buNone")) is not None


def _looks_like_body(line: str, paragraph) -> bool:
    """本文の段落か（＝箇条書きにしないか）。

    **文の終わりの記号で終わるものを本文とする。** PowerPoint は平文と
    第 1 階層の箇条書きを区別しないので、これが唯一の手掛かり。
    字下げされている段落は、書いた人が階層を意識しているので箇条書き。
    """
    return paragraph.level == 0 and line[-1] in _SENTENCE_END


def _is_mono_run(run) -> bool:
    name = (run.font.name or "").lower()
    return bool(name) and any(mono in name for mono in MONO_FONTS)


def _is_mono(shape) -> bool:
    """枠の**文字が全部**等幅か。コードブロックの手掛かり（実測: Consolas）。

    **1 つでも等幅なら、では緩すぎる。** 本文に `` `layout: 2-column` `` の
    ようなインラインコードが混ざると、その段落ごとコードブロックに
    なってしまう（実物の PowerPoint で踏んだ）。中身が全部等幅のときだけ
    コードと見なし、混ざっているものは `` ` `` で囲み直す（`_paragraph_text`）。
    """
    runs = [
        run
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    ]
    return bool(runs) and all(_is_mono_run(run) for run in runs)


def _table_blocks(table) -> list[str]:
    """表を Markdown の表にする。1 行目を見出しとして扱う。"""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [normalize_text(_cell_text(cell)).replace("\n", " ").strip() for cell in row.cells]
        rows.append([cell.replace("|", "\\|") for cell in cells])
    if not rows:
        return []

    width = len(rows[0])
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return ["\n".join(format_table(lines))]


def _cell_text(cell) -> str:
    """セルの文字。**本文と同じ書式の戻し方をする**（`` `AWS` `` を保つ）。"""
    return "\n".join(_paragraph_text(paragraph) for paragraph in cell.text_frame.paragraphs)


def _picture_blocks(shape, save_image) -> list[str]:
    """画像を保管フォルダへ渡して、本文に挿す。

    **保存先を知らないときは飛ばす。** リンクだけ書いても絵は出ない。
    """
    if save_image is None:
        return []
    try:
        image = shape.image
    except (AttributeError, KeyError, ValueError):
        logger.warning("画像を取り出せなかった")
        return []

    link = save_image(image.blob, f".{image.ext}")
    return [link] if link else []


def _notes_blocks(slide) -> list[str]:
    """発表者ノートは引用にする。**本文と混ぜない。**"""
    if not slide.has_notes_slide:
        return []
    text = normalize_text(slide.notes_slide.notes_text_frame.text).strip()
    if not text:
        return []
    return ["\n".join(f"> {line}" if line else ">" for line in text.split("\n"))]
