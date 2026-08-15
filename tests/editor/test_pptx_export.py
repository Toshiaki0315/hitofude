"""PowerPoint への書き出し（F-5）。

**ざっくり作って手で整える**前提。凝ったレイアウトは狙わない。

割り方は `core/slides.py`（F-4）が決めていて、ここは組み立てだけ。
**読み戻せること**（F-3 の取り込みを通すこと）を要件にする。往復で
形が保たれるなら、少なくとも自分の書いたものは読める。
"""

from pathlib import Path

import pytest

from hitofude.editor.pptx_export import write_pptx

pytestmark = pytest.mark.gui

DECK = """# 2026 年の計画

DX 推進チームによる提案

## プロジェクトの目的

AI を活用した業務効率の向上を目指します。

- 期間: 2026 年 1 月 〜 12 月
- 目標: 手作業の 80% 削減
    - まずは申請業務から

### 進め方

段階的に広げます。

> 最初の 3 分で目的を話す

## 実装の例

```python
def hello() -> str:
    return "こんにちは"
```

| 担当 | 人数 |
| --- | --- |
| DX | 5 名 |
"""


def png() -> bytes:
    import base64

    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
    )


@pytest.fixture
def deck_path(qapp, tmp_path: Path) -> Path:
    return write_pptx(tmp_path / "計画.pptx", DECK)


def slides(path: Path):
    from pptx import Presentation

    return list(Presentation(str(path)).slides)


def texts(slide) -> list[str]:
    return [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


class TestFile:
    def test_ファイルができる(self, deck_path) -> None:
        assert deck_path.is_file()
        assert deck_path.stat().st_size > 0

    def test_PowerPointとして開ける(self, deck_path) -> None:
        assert len(slides(deck_path)) == 3  # 表紙 + 2 枚

    def test_元の本文を変えない(self, qapp, tmp_path: Path) -> None:
        """R1: 書き出しは一方通行。"""
        text = DECK
        write_pptx(tmp_path / "a.pptx", text)
        assert text == DECK


class TestTitleSlide:
    def test_大見出しが表紙になる(self, deck_path) -> None:
        assert "2026 年の計画" in texts(slides(deck_path)[0])[0]

    def test_副題も入る(self, deck_path) -> None:
        assert any("DX 推進チーム" in text for text in texts(slides(deck_path)[0]))

    def test_大見出しが無ければ表紙を作らない(self, qapp, tmp_path: Path) -> None:
        path = write_pptx(tmp_path / "b.pptx", "## 1 枚目\n\n本文\n")
        assert len(slides(path)) == 1


class TestContent:
    def test_中見出しがスライドの題(self, deck_path) -> None:
        assert slides(deck_path)[1].shapes.title.text == "プロジェクトの目的"

    def test_段落が入る(self, deck_path) -> None:
        body = "\n".join(texts(slides(deck_path)[1]))
        assert "AI を活用した業務効率の向上を目指します。" in body

    def test_箇条書きが入る(self, deck_path) -> None:
        body = "\n".join(texts(slides(deck_path)[1]))
        assert "期間: 2026 年 1 月 〜 12 月" in body

    def test_階層が保たれる(self, deck_path) -> None:
        levels = {
            paragraph.level
            for shape in slides(deck_path)[1].shapes
            if shape.has_text_frame
            for paragraph in shape.text_frame.paragraphs
            if "申請業務" in paragraph.text
        }
        assert levels == {1}

    def test_コードが入る(self, deck_path) -> None:
        body = "\n".join(texts(slides(deck_path)[2]))
        assert 'return "こんにちは"' in body

    def test_表が入る(self, deck_path) -> None:
        tables = [shape.table for shape in slides(deck_path)[2].shapes if shape.has_table]
        assert tables
        assert tables[0].cell(0, 0).text == "担当"
        assert tables[0].cell(1, 1).text == "5 名"

    def test_発表者ノートが入る(self, deck_path) -> None:
        slide = slides(deck_path)[1]
        assert slide.has_notes_slide
        assert "最初の 3 分" in slide.notes_slide.notes_text_frame.text


class TestImages:
    def test_画像が右側に入る(self, qapp, tmp_path: Path) -> None:
        (tmp_path / "図.png").write_bytes(png())
        path = write_pptx(
            tmp_path / "図あり.pptx",
            "## 枚\n\n本文です。\n\n![](図.png)\n",
            base_path=tmp_path,
        )
        slide = slides(path)[0]
        pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
        assert pictures, "画像が入っていない"

        body = next(s for s in slide.shapes if s.has_text_frame and "本文" in s.text_frame.text)
        assert pictures[0].left > body.left, "画像が右側にない"

    def test_無い画像は飛ばす(self, qapp, tmp_path: Path) -> None:
        """**書き出しは止めない。** リンク切れでファイルが作れないのは困る。"""
        path = write_pptx(tmp_path / "欠け.pptx", "## 枚\n\n![](無い.png)\n", base_path=tmp_path)
        assert path.is_file()

    def test_保管フォルダの外は読まない(self, qapp, tmp_path: Path) -> None:
        secret = tmp_path.parent / "秘密.png"
        secret.write_bytes(png())
        base = tmp_path / "vault"
        base.mkdir()
        path = write_pptx(base / "外.pptx", "## 枚\n\n![](../秘密.png)\n", base_path=base)
        assert all(shape.shape_type != 13 for shape in slides(path)[0].shapes)


class TestRoundTrip:
    """書いたものを**自分で読み戻せる**（F-3 の取り込みを通す）。"""

    def test_題と枚数が戻る(self, deck_path) -> None:
        from hitofude.editor.pptx_import import to_markdown

        out = to_markdown(deck_path)
        assert "## プロジェクトの目的" in out
        assert "## 実装の例" in out

    def test_箇条書きが戻る(self, deck_path) -> None:
        from hitofude.editor.pptx_import import to_markdown

        assert "- 期間: 2026 年 1 月 〜 12 月" in to_markdown(deck_path)

    def test_表が戻る(self, deck_path) -> None:
        from hitofude.editor.pptx_import import to_markdown

        assert "担当" in to_markdown(deck_path)

    def test_ノートが戻る(self, deck_path) -> None:
        from hitofude.editor.pptx_import import to_markdown

        assert "> 最初の 3 分で目的を話す" in to_markdown(deck_path)


class TestEmpty:
    def test_中身が無くても壊れない(self, qapp, tmp_path: Path) -> None:
        assert write_pptx(tmp_path / "空.pptx", "").is_file()


class TestMenu:
    """書き出しの入口（`MainWindow`）。"""

    def test_メニューにある(self, qtbot, tmp_path: Path) -> None:
        from PySide6.QtCore import QSettings

        from hitofude.config import Config
        from hitofude.ui.main_window import MainWindow

        settings = QSettings(str(tmp_path / "t.ini"), QSettings.Format.IniFormat)
        config = Config(settings)
        config.vault_path = tmp_path / "Notes"
        window = MainWindow(config)
        qtbot.addWidget(window)
        try:
            assert "PowerPoint で書き出す…" in [a.text() for a in window.actions()]
        finally:
            window.close()

    def test_保管フォルダを基準に画像を探す(self, qtbot, tmp_path: Path) -> None:
        """`![](attachments/…)` は保管フォルダからの相対パス。"""
        from PySide6.QtCore import QSettings

        from hitofude.config import Config
        from hitofude.ui.main_window import MainWindow

        settings = QSettings(str(tmp_path / "t.ini"), QSettings.Format.IniFormat)
        config = Config(settings)
        config.vault_path = tmp_path / "Notes"
        window = MainWindow(config)
        qtbot.addWidget(window)
        try:
            (window.vault.attachments_dir / "図.png").write_bytes(png())
            target = window._write_pptx(
                tmp_path / "出力.pptx", "## 枚\n\n![](attachments/図.png)\n"
            )
            assert any(shape.shape_type == 13 for shape in slides(target)[0].shapes)
        finally:
            window.close()
