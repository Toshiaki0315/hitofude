"""PowerPoint の読み込み（F-3）。

**ざっくり読んで手で直す**前提。元のレイアウトは復元しない。

手掛かりは実物（md2pptx が作った `.pptx`）を調べて決めた。

| 手掛かり | 使い道 |
| --- | --- |
| スライドのタイトル枠 | `##` |
| `buNone`（行頭記号なし）で短い段落 | `###` |
| 文の終わりの記号で終わる段落 | 本文。それ以外は `- ` |
| 等幅フォント（Consolas など） | コードブロック |
| 太字の run | `**強調**` |

**他のツールで作った `.pptx` でも落ちない**ことを要件にする。手掛かりが
無ければ、ただの段落として拾えばよい。
"""

from pathlib import Path

import pytest

from hitofude.editor.pptx_import import to_markdown

pytestmark = pytest.mark.gui


def build(path: Path, *, notes: str = "") -> Path:
    """試験用の `.pptx` を組む。**実物の形に寄せる。**"""
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "現状の分析と課題"

    frame = slide.placeholders[1].text_frame
    frame.text = "現場では多くの課題が山積しています。"

    heading = frame.add_paragraph()
    heading.text = "システムの老朽化"
    heading.level = 0
    _drop_bullet(heading)

    item = frame.add_paragraph()
    item.text = "既存システムの動作が遅い"

    nested = frame.add_paragraph()
    nested.text = "担当者が不在"
    nested.level = 1

    box = slide.shapes.add_textbox(Pt(10), Pt(300), Pt(400), Pt(100))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = 'def hello():\n    return "こんにちは"'
    run.font.name = "Consolas"

    number = slide.shapes.add_textbox(Pt(10), Pt(500), Pt(40), Pt(20))
    number.text_frame.text = "1"

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    presentation.save(str(path))
    return path


def _drop_bullet(paragraph) -> None:
    """行頭記号を消す（`buNone`）。md2pptx が小見出しに使う印。"""
    from pptx.oxml.ns import qn

    properties = paragraph._pPr if paragraph._pPr is not None else paragraph._p.get_or_add_pPr()
    properties.append(properties.makeelement(qn("a:buNone"), {}))


@pytest.fixture
def sample(qapp, tmp_path: Path) -> Path:
    return build(tmp_path / "講演資料.pptx")


class TestStructure:
    def test_題名はファイル名(self, sample) -> None:
        assert to_markdown(sample).startswith("# 講演資料\n")

    def test_スライドは中見出し(self, sample) -> None:
        assert "## 現状の分析と課題" in to_markdown(sample)

    def test_行頭記号の無い短い段落は小見出し(self, sample) -> None:
        """md2pptx は `###` を `buNone` で書く。"""
        assert "### システムの老朽化" in to_markdown(sample)

    def test_文で終わる段落は本文(self, sample) -> None:
        out = to_markdown(sample)
        assert "\n現場では多くの課題が山積しています。\n" in out

    def test_文で終わらない段落は箇条書き(self, sample) -> None:
        assert "- 既存システムの動作が遅い" in to_markdown(sample)

    def test_階層は字下げする(self, sample) -> None:
        assert "    - 担当者が不在" in to_markdown(sample)

    def test_ページ番号は入らない(self, sample) -> None:
        out = to_markdown(sample)
        assert "\n1\n" not in out
        assert "- 1" not in out


class TestCode:
    def test_等幅の枠はコードになる(self, sample) -> None:
        out = to_markdown(sample)
        assert "```\ndef hello():" in out

    def test_中身はそのまま(self, sample) -> None:
        """**コードの中は触らない。** 箇条書きにも段落にもしない。"""
        assert '    return "こんにちは"\n```' in to_markdown(sample)


class TestInlineCode:
    """インラインコードが混ざっても、段落はコードブロックにしない。

    実物で踏んだ。本文に `` `layout: 2-column` `` があるだけで、その枠が
    丸ごとコードブロックになっていた（**1 つでも等幅なら**、という判定が
    緩すぎた）。
    """

    def build_mixed(self, tmp_path: Path) -> Path:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "設定"
        paragraph = slide.placeholders[1].text_frame.paragraphs[0]
        head = paragraph.add_run()
        head.text = "明示的に "
        code = paragraph.add_run()
        code.text = "layout: 2-column"
        code.font.name = "Consolas"
        tail = paragraph.add_run()
        tail.text = " を指定します。"

        path = tmp_path / "混在.pptx"
        presentation.save(str(path))
        return path

    def test_コードブロックにしない(self, qapp, tmp_path: Path) -> None:
        assert "```" not in to_markdown(self.build_mixed(tmp_path))

    def test_記号で囲み直す(self, qapp, tmp_path: Path) -> None:
        out = to_markdown(self.build_mixed(tmp_path))
        assert "明示的に `layout: 2-column` を指定します。" in out


class TestNotes:
    def test_発表者ノートは引用になる(self, qapp, tmp_path: Path) -> None:
        path = build(tmp_path / "ノート付き.pptx", notes="最初の 3 分で目的を話す")
        assert "> 最初の 3 分で目的を話す" in to_markdown(path)

    def test_ノートが無ければ何も出ない(self, sample) -> None:
        assert ">" not in to_markdown(sample)


class TestTable:
    def test_表になる(self, qapp, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Pt

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "役割分担"
        table = slide.shapes.add_table(2, 2, Pt(10), Pt(100), Pt(400), Pt(100)).table
        table.cell(0, 0).text = "担当"
        table.cell(0, 1).text = "人数"
        table.cell(1, 0).text = "DX 推進"
        table.cell(1, 1).text = "5 名"
        path = tmp_path / "表.pptx"
        presentation.save(str(path))

        out = to_markdown(path)
        assert "| 担当" in out
        assert "| --- |" in out.replace("-" * 4, "---").replace("--- |", "--- |")
        assert "DX 推進" in out


class TestImages:
    def png(self, tmp_path: Path) -> Path:
        """**1×1 では試せない。** 飾りとして間引かれる（ADR-0027 追記）ので、
        取り込む価値のある大きさで作る。"""
        from PIL import Image

        path = tmp_path / "素材.png"
        Image.new("RGB", (200, 200), (30, 90, 180)).save(path)
        return path

    def build_with_image(self, tmp_path: Path) -> Path:
        from pptx import Presentation
        from pptx.util import Pt

        image = tmp_path / "図.png"
        image.write_bytes(self.png(tmp_path).read_bytes())
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "構成図"
        slide.shapes.add_picture(str(image), Pt(10), Pt(100), Pt(200), Pt(150))
        path = tmp_path / "図あり.pptx"
        presentation.save(str(path))
        return path

    def test_保存先を渡せば画像が入る(self, qapp, tmp_path: Path) -> None:
        saved: list[tuple[bytes, str]] = []

        def keep(data: bytes, suffix: str) -> str:
            saved.append((data, suffix))
            return "![](attachments/図.png)"

        out = to_markdown(self.build_with_image(tmp_path), save_image=keep)
        assert saved and saved[0][1] == ".png"
        assert "![](attachments/図.png)" in out

    def test_保存先が無ければ画像は飛ばす(self, qapp, tmp_path: Path) -> None:
        """**落とさない。** 画像を置く場所を知らないだけ。"""
        out = to_markdown(self.build_with_image(tmp_path))
        assert "構成図" in out
        assert "![](" not in out

    def test_保存に失敗しても止まらない(self, qapp, tmp_path: Path) -> None:
        out = to_markdown(self.build_with_image(tmp_path), save_image=lambda data, suffix: None)
        assert "構成図" in out


class TestBroken:
    def test_PPTXでないファイルは空(self, qapp, tmp_path: Path) -> None:
        broken = tmp_path / "偽物.pptx"
        broken.write_text("これは PowerPoint ではありません", encoding="utf-8")
        assert to_markdown(broken) == ""

    def test_無いファイルは空(self, qapp, tmp_path: Path) -> None:
        assert to_markdown(tmp_path / "無い.pptx") == ""

    def test_中身の無いスライドだけなら空(self, qapp, tmp_path: Path) -> None:
        """画像だけの PDF と同じ扱い（題名だけのノートを作らせない）。"""
        from pptx import Presentation

        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        path = tmp_path / "空.pptx"
        presentation.save(str(path))
        assert to_markdown(path) == ""


class TestPptxImages:
    """PowerPoint の画像も PDF と同じ規則で間引く（ユーザー指摘 2026-08-23）。

    **PDF だけ間引いていた。** 同じ絵が 2 枚入り、40px の飾りまで取り込まれて
    いた（実測）。取り込みの入口で作法が違うのは、後から読む人が混乱する。
    """

    def deck(self, tmp_path, pictures):
        """`pictures` は `(幅 px, 中身の種)` の並び。"""
        import hashlib

        from PIL import Image
        from pptx import Presentation
        from pptx.util import Cm

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "図のスライド"
        for index, (width, seed) in enumerate(pictures):
            source = tmp_path / f"{seed}-{width}.png"
            if not source.exists():
                # **種が同じなら同じ絵、違えば違う絵**（色を種から作る）。
                # `hash()` を 200 で割ると 35 枚では**必ず衝突する**（誕生日
                # 問題）。別のはずの絵が同じ中身になり、間引きに落ちて
                # テストが揺れた。1600 万色に散らす
                tone = hashlib.md5(seed.encode()).digest()[:3]
                Image.new("RGB", (width, width), tuple(tone)).save(source)
            slide.shapes.add_picture(str(source), Cm(1 + index * 5), Cm(5), Cm(3))
        path = tmp_path / "図つき.pptx"
        presentation.save(str(path))
        return path

    def saved_from(self, path):
        from hitofude.editor import importer

        saved: list[bytes] = []
        importer.to_markdown(path, save_image=lambda data, _s: (saved.append(data), "![](x)")[1])
        return saved

    def test_同じ絵は一度だけ(self, tmp_path) -> None:
        path = self.deck(tmp_path, [(400, "同じ"), (400, "同じ"), (400, "別")])
        assert len(self.saved_from(path)) == 2

    def test_小さい絵は入らない(self, tmp_path) -> None:
        from hitofude.core.imported import MIN_IMAGE_SIDE

        path = self.deck(tmp_path, [(MIN_IMAGE_SIDE - 40, "飾り"), (400, "図")])
        assert len(self.saved_from(path)) == 1

    def test_多すぎれば打ち切る(self, tmp_path) -> None:
        from hitofude.core.imported import MAX_IMAGES

        path = self.deck(tmp_path, [(400, f"図{n}") for n in range(MAX_IMAGES + 5)])
        assert len(self.saved_from(path)) == MAX_IMAGES
